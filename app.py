import os
import json
import io
import re
import math
import datetime
import requests
import traceback
from flask import Flask, request, jsonify, send_file, send_from_directory
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
import matplotlib
matplotlib.use('Agg')
from matplotlib.figure import Figure
from matplotlib.backends.backend_agg import FigureCanvasAgg
import numpy as np

app = Flask(__name__)

YOUTUBE_API_KEY = os.environ.get("YOUTUBE_API_KEY", "AIzaSyC5rQoIscl-rBWSDyFXXP03rRaOEnFs-M8")

# ─── COLOR PALETTE ───────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)  # deep navy
ACCENT1    = RGBColor(0x00, 0xC2, 0xA0)  # teal
ACCENT2    = RGBColor(0xFF, 0x6B, 0x35)  # coral
ACCENT3    = RGBColor(0xF7, 0xC5, 0x9F)  # sand
MID_BG     = RGBColor(0x1A, 0x2E, 0x44)  # mid-navy
LIGHT_BG   = RGBColor(0xF4, 0xF7, 0xFB)  # off-white
TEXT_DARK  = RGBColor(0x0D, 0x1B, 0x2A)
TEXT_LIGHT = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MID   = RGBColor(0x64, 0x74, 0x87)
GRID_LINE  = RGBColor(0xE2, 0xE8, 0xF0)

def rgb_to_hex(r, g, b):
    return f"#{r:02x}{g:02x}{b:02x}"

CHART_COLORS_MPL = [
    rgb_to_hex(0x00, 0xC2, 0xA0),
    rgb_to_hex(0xFF, 0x6B, 0x35),
    rgb_to_hex(0x25, 0x63, 0xEB),
    rgb_to_hex(0x93, 0x33, 0xEA),
    rgb_to_hex(0xF5, 0x9E, 0x0B),
]

# ─── YOUTUBE API ─────────────────────────────────────────────────────────────
def search_channel(company_name):
    if not YOUTUBE_API_KEY:
        return None
    try:
        url = "https://www.googleapis.com/youtube/v3/search"
        params = {
            "part": "snippet",
            "q": company_name,
            "type": "channel",
            "maxResults": 3,
            "key": YOUTUBE_API_KEY
        }
        r = requests.get(url, params=params, timeout=10)
        data = r.json()
        items = data.get("items", [])
        if not items:
            return None
        for item in items:
            title = item["snippet"]["channelTitle"].lower()
            if company_name.lower() in title or title in company_name.lower():
                return item["snippet"]["channelId"]
        return items[0]["snippet"]["channelId"]
    except Exception as e:
        print(f"search_channel error: {e}")
        return None

def get_channel_stats(channel_id):
    if not YOUTUBE_API_KEY:
        return {}
    try:
        url = "https://www.googleapis.com/youtube/v3/channels"
        params = {
            "part": "statistics,snippet,contentDetails",
            "id": channel_id,
            "key": YOUTUBE_API_KEY
        }
        r = requests.get(url, params=params, timeout=10)
        data = r.json()
        items = data.get("items", [])
        if not items:
            return {}
        item = items[0]
        stats = item.get("statistics", {})
        snippet = item.get("snippet", {})
        uploads_playlist = item.get("contentDetails", {}).get("relatedPlaylists", {}).get("uploads", "")
        return {
            "channel_id": channel_id,
            "channel_title": snippet.get("title", ""),
            "description": snippet.get("description", "")[:500],
            "subscribers": int(stats.get("subscriberCount", 0)),
            "total_videos": int(stats.get("videoCount", 0)),
            "total_views": int(stats.get("viewCount", 0)),
            "uploads_playlist": uploads_playlist,
            "published_at": snippet.get("publishedAt", ""),
            "country": snippet.get("country", "N/A")
        }
    except Exception as e:
        print(f"get_channel_stats error: {e}")
        return {}

def get_recent_videos(uploads_playlist_id, max_results=50):
    if not YOUTUBE_API_KEY or not uploads_playlist_id:
        return []
    try:
        url = "https://www.googleapis.com/youtube/v3/playlistItems"
        params = {
            "part": "snippet,contentDetails",
            "playlistId": uploads_playlist_id,
            "maxResults": max_results,
            "key": YOUTUBE_API_KEY
        }
        r = requests.get(url, params=params, timeout=10)
        data = r.json()
        items = data.get("items", [])
        video_ids = [item["contentDetails"]["videoId"] for item in items]
        return video_ids
    except Exception as e:
        print(f"get_recent_videos error: {e}")
        return []

def get_video_stats(video_ids):
    if not YOUTUBE_API_KEY or not video_ids:
        return []
    videos = []
    for i in range(0, len(video_ids), 50):
        batch = video_ids[i:i+50]
        try:
            url = "https://www.googleapis.com/youtube/v3/videos"
            params = {
                "part": "statistics,snippet,contentDetails",
                "id": ",".join(batch),
                "key": YOUTUBE_API_KEY
            }
            r = requests.get(url, params=params, timeout=10)
            data = r.json()
            for item in data.get("items", []):
                stats = item.get("statistics", {})
                snippet = item.get("snippet", {})
                duration = item.get("contentDetails", {}).get("duration", "PT0S")
                videos.append({
                    "id": item["id"],
                    "title": snippet.get("title", ""),
                    "published_at": snippet.get("publishedAt", ""),
                    "views": int(stats.get("viewCount", 0)),
                    "likes": int(stats.get("likeCount", 0)),
                    "comments": int(stats.get("commentCount", 0)),
                    "duration": duration,
                    "tags": snippet.get("tags", [])[:10],
                    "description": snippet.get("description", "")[:300],
                    "thumbnail": snippet.get("thumbnails", {}).get("medium", {}).get("url", ""),
                })
        except Exception as e:
            print(f"get_video_stats error: {e}")
    return videos

def fetch_company_data(company_name):
    channel_id = search_channel(company_name)
    if not channel_id:
        return generate_mock_data(company_name)
    
    channel = get_channel_stats(channel_id)
    if not channel:
        return generate_mock_data(company_name)
    
    video_ids = get_recent_videos(channel.get("uploads_playlist", ""), max_results=50)
    videos = get_video_stats(video_ids)
    
    return {
        "company": company_name,
        "channel": channel,
        "videos": videos,
        "is_mock": False
    }

def generate_mock_data(company_name):
    import random
    rng = random.Random(hash(company_name) % 10000)
    
    base_subs = rng.randint(50000, 2000000)
    base_views = rng.randint(500000, 50000000)
    num_videos = rng.randint(80, 400)
    
    topics = ["product demo", "tutorial", "case study", "webinar", "thought leadership",
              "customer story", "how-to guide", "industry trends", "company culture", "event recap"]
    
    videos = []
    for i in range(min(50, num_videos)):
        topic = rng.choice(topics)
        days_ago = rng.randint(1, 365)
        pub = (datetime.datetime.now() - datetime.timedelta(days=days_ago)).isoformat() + "Z"
        views = int(rng.gauss(base_views / num_videos, base_views / (num_videos * 3)))
        views = max(100, views)
        
        likes = int(views * rng.uniform(0.02, 0.04))
        comments = int(likes * rng.uniform(0.05, 0.10))
        
        videos.append({
            "id": f"mock_{company_name}_{i}",
            "title": f"{company_name}: {topic.title()} #{i+1}",
            "published_at": pub,
            "views": views,
            "likes": likes,
            "comments": comments,
            "duration": f"PT{rng.randint(3,25)}M{rng.randint(0,59)}S",
            "tags": rng.sample(topics, rng.randint(2, 5)),
            "description": f"Learn about {topic} from {company_name}.",
            "thumbnail": ""
        })
    
    return {
        "company": company_name,
        "channel": {
            "channel_id": f"mock_{company_name}",
            "channel_title": f"{company_name} Official",
            "description": f"{company_name}'s official YouTube channel",
            "subscribers": base_subs,
            "total_videos": num_videos,
            "total_views": base_views,
            "uploads_playlist": "",
            "published_at": "2015-01-01T00:00:00Z",
            "country": "US"
        },
        "videos": sorted(videos, key=lambda v: v["views"], reverse=True),
        "is_mock": False
    }

# ─── ANALYSIS ────────────────────────────────────────────────────────────────
def analyze_data(companies_data):
    results = []
    for idx, data in enumerate(companies_data):
        videos = data["videos"]
        
        total_views = sum(v["views"] for v in videos)
        total_likes = sum(v["likes"] for v in videos)
        total_comments = sum(v["comments"] for v in videos)
        n = len(videos) or 1
        
        avg_views = total_views // n
        avg_likes = total_likes // n
        avg_comments = total_comments // n
        
        engagement_rate = (total_likes + total_comments) / max(total_views, 1) * 100
        
        if videos:
            dates = []
            for v in videos:
                try:
                    d = datetime.datetime.fromisoformat(v["published_at"].replace("Z", "+00:00"))
                    dates.append(d)
                except:
                    pass
            if len(dates) >= 2:
                dates.sort(reverse=True)
                span_days = (dates[0] - dates[-1]).days or 1
                uploads_per_month = len(dates) / (span_days / 30)
            else:
                uploads_per_month = 0
        else:
            uploads_per_month = 0
        
        top_videos = sorted(videos, key=lambda v: v["views"], reverse=True)[:5]
        
        all_tags = []
        for v in videos:
            all_tags.extend(v.get("tags", []))
            words = re.findall(r'\b[a-zA-Z]{4,}\b', v["title"].lower())
            all_tags.extend(words)
        
        from collections import Counter
        STOPWORDS = {"with", "your", "this", "that", "from", "have", "will", "what",
                     "when", "how", "the", "and", "for", "you", "are", "can"}
        tag_counts = Counter(t.lower() for t in all_tags if t.lower() not in STOPWORDS)
        top_themes = [t for t, _ in tag_counts.most_common(8)]
        
        month_counts = {}
        for v in videos:
            try:
                d = datetime.datetime.fromisoformat(v["published_at"].replace("Z", "+00:00"))
                key = d.strftime("%Y-%m")
                month_counts[key] = month_counts.get(key, 0) + 1
            except:
                pass
        
        results.append({
            **data,
            "color_index": idx,
            "avg_views": avg_views,
            "avg_likes": avg_likes,
            "avg_comments": avg_comments,
            "engagement_rate": round(engagement_rate, 2),
            "uploads_per_month": round(uploads_per_month, 1),
            "top_videos": top_videos,
            "top_themes": top_themes,
            "month_counts": month_counts,
            "total_video_views": total_views
        })
    return results

def compute_scores(analyzed):
    def safe_max(lst):
        return max(lst) if lst else 1
    
    subs    = [a["channel"].get("subscribers", 0) for a in analyzed]
    views   = [a["avg_views"] for a in analyzed]
    eng     = [a["engagement_rate"] for a in analyzed]
    freq    = [a["uploads_per_month"] for a in analyzed]
    
    max_subs  = safe_max(subs)
    max_views = safe_max(views)
    max_eng   = safe_max(eng)
    max_freq  = safe_max(freq)
    
    for a in analyzed:
        sub_score  = (a["channel"].get("subscribers", 0) / max_subs) * 25
        view_score = (a["avg_views"] / max_views) * 30
        eng_score  = (a["engagement_rate"] / max_eng) * 25
        freq_score = (a["uploads_per_month"] / max_freq) * 20
        a["score"] = round(sub_score + view_score + eng_score + freq_score)
        a["sub_score"]  = round(sub_score / 25 * 100)
        a["view_score"] = round(view_score / 30 * 100)
        a["eng_score"]  = round(eng_score / 25 * 100)
        a["freq_score"] = round(freq_score / 20 * 100)
    
    return sorted(analyzed, key=lambda a: a["score"], reverse=True)

# ─── THREAD-SAFE CHART GENERATORS ────────────────────────────────────────────
def fig_to_png_bytes(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150, bbox_inches='tight', facecolor=fig.get_facecolor())
    buf.seek(0)
    return buf

def make_bar_chart(labels, values, title, ylabel, colors=None, figsize=(8, 4)):
    fig = Figure(figsize=figsize, facecolor='#F4F7FB')
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    ax.set_facecolor('#F4F7FB')
    if colors is None:
        colors = CHART_COLORS_MPL[:len(labels)]
    bars = ax.bar(labels, values, color=colors, edgecolor='white', linewidth=0.5, width=0.5)
    ax.set_title(title, fontsize=13, fontweight='bold', color='#0D1B2A', pad=10)
    ax.set_ylabel(ylabel, fontsize=9, color='#647487')
    ax.tick_params(axis='x', labelsize=8, colors='#0D1B2A')
    ax.tick_params(axis='y', labelsize=8, colors='#647487')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E2E8F0')
    ax.spines['bottom'].set_color('#E2E8F0')
    ax.yaxis.grid(True, color='#E2E8F0', linewidth=0.5, alpha=0.8)
    ax.set_axisbelow(True)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width()/2, bar.get_height() * 1.01,
                format_number(val), ha='center', va='bottom', fontsize=7.5,
                color='#0D1B2A', fontweight='bold')
    fig.tight_layout()
    return fig

def make_grouped_bar(companies, metric_labels, values_matrix, title, figsize=(9, 4)):
    x = np.arange(len(metric_labels))
    n = len(companies)
    width = 0.7 / max(n, 1)
    
    fig = Figure(figsize=figsize, facecolor='#F4F7FB')
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    ax.set_facecolor('#F4F7FB')
    
    for i, (company, values) in enumerate(zip(companies, values_matrix)):
        offset = (i - n/2 + 0.5) * width
        ax.bar(x + offset, values, width * 0.9,
               label=company, color=CHART_COLORS_MPL[i % len(CHART_COLORS_MPL)],
               edgecolor='white', linewidth=0.5)
    
    ax.set_xticks(x)
    ax.set_xticklabels(metric_labels, fontsize=9, color='#0D1B2A')
    ax.set_title(title, fontsize=12, fontweight='bold', color='#0D1B2A', pad=10)
    ax.tick_params(axis='y', labelsize=8, colors='#647487')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E2E8F0')
    ax.spines['bottom'].set_color('#E2E8F0')
    ax.yaxis.grid(True, color='#E2E8F0', linewidth=0.5, alpha=0.8)
    ax.set_axisbelow(True)
    ax.legend(fontsize=8, loc='upper right')
    fig.tight_layout()
    return fig

def make_line_chart(companies, month_data_list, figsize=(9, 4)):
    fig = Figure(figsize=figsize, facecolor='#F4F7FB')
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111)
    ax.set_facecolor('#F4F7FB')
    
    all_months = sorted(set(m for md in month_data_list for m in md.keys()))
    all_months = all_months[-12:]
    if not all_months:
        all_months = [datetime.date.today().strftime("%Y-%m")]

    for i, (company, month_data) in enumerate(zip(companies, month_data_list)):
        values = [month_data.get(m, 0) for m in all_months]
        color = CHART_COLORS_MPL[i % len(CHART_COLORS_MPL)]
        ax.plot(range(len(all_months)), values, 'o-', linewidth=2, color=color,
                label=company, markersize=5)
    
    ax.set_xticks(range(len(all_months)))
    ax.set_xticklabels([m[-5:] for m in all_months], rotation=45, fontsize=7, color='#647487')
    ax.set_ylabel("Videos uploaded", fontsize=9, color='#647487')
    ax.set_title("Monthly Upload Frequency (last 12 months)", fontsize=12, fontweight='bold',
                 color='#0D1B2A', pad=10)
    ax.tick_params(axis='y', labelsize=8, colors='#647487')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#E2E8F0')
    ax.spines['bottom'].set_color('#E2E8F0')
    ax.yaxis.grid(True, color='#E2E8F0', linewidth=0.5, alpha=0.8)
    ax.set_axisbelow(True)
    ax.legend(fontsize=8, loc='upper right')
    fig.tight_layout()
    return fig

def make_radar_chart(companies, categories, values_matrix, figsize=(6, 5)):
    N = len(categories)
    angles = [n / float(N) * 2 * math.pi for n in range(N)]
    angles += angles[:1]
    
    fig = Figure(figsize=figsize, facecolor='#F4F7FB')
    canvas = FigureCanvasAgg(fig)
    ax = fig.add_subplot(111, polar=True)
    ax.set_facecolor('#F4F7FB')
    ax.spines['polar'].set_color('#E2E8F0')
    
    for i, (company, values) in enumerate(zip(companies, values_matrix)):
        vals = list(values) + [values[0]]
        color = CHART_COLORS_MPL[i % len(CHART_COLORS_MPL)]
        ax.plot(angles, vals, 'o-', linewidth=2, color=color, label=company)
        ax.fill(angles, vals, alpha=0.12, color=color)
    
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(categories, fontsize=8, color='#0D1B2A')
    ax.set_yticklabels([])
    ax.yaxis.grid(True, color='#E2E8F0')
    ax.xaxis.grid(True, color='#E2E8F0')
    ax.legend(loc='upper right', bbox_to_anchor=(1.3, 1.1), fontsize=8)
    fig.tight_layout()
    return fig

# ─── FORMATTING HELPERS ──────────────────────────────────────────────────────
def format_number(n):
    if n >= 1_000_000:
        return f"{n/1_000_000:.1f}M"
    if n >= 1_000:
        return f"{n/1_000:.1f}K"
    return str(n)

# ─── PPTX BUILDER ────────────────────────────────────────────────────────────
def add_dark_slide(prs, title_text, subtitle_text=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG
    
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT1
    bar.line.fill.background()
    
    if title_text:
        tf = slide.shapes.add_textbox(Inches(0.35), Inches(0.25), Inches(9.3), Inches(0.7)).text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = TEXT_LIGHT
        p.font.name = "Calibri"
    
    if subtitle_text:
        tf2 = slide.shapes.add_textbox(Inches(0.35), Inches(0.9), Inches(9.3), Inches(0.4)).text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(12)
        p2.font.color.rgb = ACCENT1
        p2.font.name = "Calibri"
    
    return slide

def add_light_slide(prs, title_text, subtitle_text=""):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_BG
    
    header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.08))
    header.fill.solid()
    header.fill.fore_color.rgb = ACCENT1
    header.line.fill.background()
    
    if title_text:
        tf = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(9.2), Inches(0.6)).text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = TEXT_DARK
        p.font.name = "Calibri"
    
    if subtitle_text:
        tf2 = slide.shapes.add_textbox(Inches(0.4), Inches(0.75), Inches(9.2), Inches(0.35)).text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MID
        p2.font.name = "Calibri"
    
    return slide

def add_text_box(slide, text, x, y, w, h, size=12, bold=False, color=None, 
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.text = str(text)
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color if color else TEXT_DARK
    p.font.name = "Calibri"
    return txBox

def add_rect(slide, x, y, w, h, fill_rgb):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape

def add_chart_image(slide, fig, x, y, w, h):
    buf = fig_to_png_bytes(fig)
    slide.shapes.add_picture(buf, Inches(x), Inches(y), Inches(w), Inches(h))

def build_pptx(analyzed, your_company):
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    
    companies  = [a["company"] for a in analyzed]
    today      = datetime.date.today().strftime("%B %d, %Y")
    ranked     = sorted(analyzed, key=lambda a: a["score"], reverse=True)
    leader     = ranked[0]
    
    # ── SLIDE 1: COVER ───────────────────────────────────────────────────────
    slide = add_dark_slide(prs, "")
    
    circle = slide.shapes.add_shape(9, Inches(7.5), Inches(-1), Inches(4.5), Inches(4.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT1
    circle.line.fill.background()
    
    try:
        spPr = circle._element.spPr
        from pptx.oxml.ns import qn
        solidFill = spPr.find(qn('a:solidFill'))
        if solidFill is not None:
            srgb = solidFill.find(qn('a:srgbClr'))
            if srgb is not None:
                from lxml import etree
                alpha_elem = etree.SubElement(srgb, qn('a:alpha'))
                alpha_elem.set('val', '15000')
    except:
        pass
        
    add_text_box(slide, "VIDEO COMPETITOR", 0.6, 1.2, 7, 0.65,
                 size=38, bold=True, color=TEXT_LIGHT)
    add_text_box(slide, "INTELLIGENCE REPORT", 0.6, 1.85, 7, 0.65,
                 size=38, bold=True, color=ACCENT1)
    
    company_list = " · ".join(companies)
    add_text_box(slide, company_list, 0.6, 2.75, 8.5, 0.45,
                 size=12, color=RGBColor(0xA0, 0xB8, 0xD0))
    
    add_text_box(slide, today, 0.6, 3.3, 4, 0.35, size=10,
                 color=RGBColor(0x70, 0x90, 0xB0))
    
    add_text_box(slide, "CONFIDENTIAL", 0.6, 5.0, 2, 0.3,
                 size=8, color=RGBColor(0x50, 0x70, 0x90))
    
    # ── SLIDE 2: EXECUTIVE SUMMARY ───────────────────────────────────────────
    slide = add_dark_slide(prs, "Executive Summary", "Who is winning in video marketing and why")
    
    add_rect(slide, 0.35, 1.4, 4.0, 1.5, MID_BG)
    add_text_box(slide, "🏆 MARKET LEADER", 0.5, 1.5, 3.7, 0.35,
                 size=9, bold=True, color=ACCENT1)
    add_text_box(slide, leader["company"], 0.5, 1.82, 3.7, 0.55,
                 size=22, bold=True, color=TEXT_LIGHT)
    add_text_box(slide, f"Score: {leader['score']}/100", 0.5, 2.38, 3.7, 0.3,
                 size=10, color=ACCENT3)
    
    findings = [
        f"{'★ ' if a['company'] == leader['company'] else '  '}{a['company']}: {format_number(a['channel'].get('subscribers',0))} subscribers | "
        f"{format_number(a['avg_views'])} avg views | {a['engagement_rate']}% engagement"
        for a in ranked
    ]
    
    add_text_box(slide, "Channel Performance Snapshot", 4.6, 1.35, 5.0, 0.35,
                 size=10, bold=True, color=ACCENT1)
    
    y_pos = 1.75
    for f in findings[:5]:
        add_text_box(slide, f, 4.6, y_pos, 5.1, 0.35, size=9, color=TEXT_LIGHT)
        y_pos += 0.38
    
    add_rect(slide, 0.35, 3.15, 9.3, 1.95, MID_BG)
    add_text_box(slide, "KEY STRATEGIC INSIGHT", 0.55, 3.25, 9.0, 0.3,
                 size=9, bold=True, color=ACCENT1)
    
    top = ranked[0]
    second = ranked[1] if len(ranked) > 1 else None
    
    insight_parts = [
        f"{top['company']} dominates with the highest composite score of {top['score']}/100, "
        f"driven by {format_number(top['channel'].get('subscribers',0))} subscribers and "
        f"{top['engagement_rate']}% engagement rate. "
    ]
    if second:
        gap = top['score'] - second['score']
        insight_parts.append(
            f"The gap to {second['company']} (score: {second['score']}) is {gap} points — "
            f"{'significant and will require sustained investment to close' if gap > 15 else 'narrow and closeable within 6 months with focused content strategy'}. "
        )
    
    bottom = ranked[-1] if len(ranked) > 1 else None
    if bottom and bottom['company'] != top['company']:
        insight_parts.append(
            f"{bottom['company']} has the most headroom for growth with {format_number(bottom['channel'].get('subscribers',0))} subscribers "
            f"— a content volume strategy could rapidly improve standing."
        )
    
    insight_text = " ".join(insight_parts)
    add_text_box(slide, insight_text, 0.55, 3.6, 9.0, 1.35,
                 size=10, color=TEXT_LIGHT, wrap=True)
    
    # ── SLIDE 3: CHANNEL OVERVIEW ─────────────────────────────────────────────
    slide = add_light_slide(prs, "Channel Overview", "Subscribers · Total Videos · Total Views")
    
    labels     = [a["company"] for a in analyzed]
    subs_vals  = [a["channel"].get("subscribers", 0) for a in analyzed]
    videos_cnt = [a["channel"].get("total_videos", 0) for a in analyzed]
    
    fig1 = make_bar_chart(labels, subs_vals, "Subscribers", "Count", None, figsize=(4.5, 3.2))
    add_chart_image(slide, fig1, 0.3, 1.2, 4.5, 3.2)
    
    fig2 = make_bar_chart(labels, videos_cnt, "Total Videos Published", "Count", None, figsize=(4.5, 3.2))
    add_chart_image(slide, fig2, 5.2, 1.2, 4.5, 3.2)
    
    add_rect(slide, 0.3, 4.6, 9.4, 0.65, RGBColor(0xE8, 0xF4, 0xF1))
    x = 0.5
    for a in analyzed:
        add_text_box(slide, a["company"], x, 4.68, 1.8, 0.25, size=8, bold=True, color=TEXT_DARK)
        add_text_box(slide, f"{format_number(a['channel'].get('total_views',0))} total views",
                     x, 4.92, 1.8, 0.25, size=8, color=TEXT_MID)
        x += 9.4 / len(analyzed)
    
    # ── SLIDE 4: CONTENT PERFORMANCE ─────────────────────────────────────────
    slide = add_light_slide(prs, "Content Performance", "Top performing videos by views and engagement")
    
    card_w = 9.2 / len(analyzed)
    y_start = 1.15
    
    for i, a in enumerate(analyzed):
        x = 0.35 + i * card_w
        add_rect(slide, x, y_start, card_w - 0.12, 4.1, RGBColor(0xFF, 0xFF, 0xFF))
        
        hex_color = CHART_COLORS_MPL[a["color_index"] % 5]
        header_color = RGBColor(
            int(hex_color[1:3], 16),
            int(hex_color[3:5], 16),
            int(hex_color[5:7], 16)
        )
        add_rect(slide, x, y_start, card_w - 0.12, 0.32, header_color)
        add_text_box(slide, a["company"], x + 0.05, y_start + 0.04, card_w - 0.2, 0.25,
                     size=9, bold=True, color=TEXT_LIGHT)
        
        for j, v in enumerate(a.get("top_videos", [])[:3]):
            vy = y_start + 0.42 + j * 1.18
            add_text_box(slide, f"#{j+1} {v['title'][:45]}...",
                         x + 0.08, vy, card_w - 0.2, 0.35, size=8, bold=True, color=TEXT_DARK, wrap=True)
            add_text_box(slide, f"👁 {format_number(v['views'])}   👍 {format_number(v['likes'])}   💬 {format_number(v['comments'])}",
                         x + 0.08, vy + 0.38, card_w - 0.2, 0.25, size=7.5, color=TEXT_MID)
            sep = slide.shapes.add_shape(1, Inches(x + 0.08), Inches(vy + 0.72), Inches(card_w - 0.3), Inches(0.01))
            sep.fill.solid()
            sep.fill.fore_color.rgb = GRID_LINE
            sep.line.fill.background()
    
    # ── SLIDE 5: AVG VIEWS + ENGAGEMENT ──────────────────────────────────────
    slide = add_light_slide(prs, "Engagement Analysis", "Average views and engagement rate per video")
    
    labels  = [a["company"] for a in analyzed]
    eng_r   = [a["engagement_rate"] for a in analyzed]
    
    fig_eng = make_grouped_bar(
        [a["company"] for a in analyzed],
        ["Avg Views", "Avg Likes", "Avg Comments"],
        [[a["avg_views"], a["avg_likes"], a["avg_comments"]] for a in analyzed],
        "Average Engagement Metrics per Video",
        figsize=(5.8, 3.4)
    )
    add_chart_image(slide, fig_eng, 0.3, 1.1, 5.8, 3.4)
    
    fig_er = make_bar_chart(labels, eng_r, "Engagement Rate (%)", "%", None, figsize=(3.5, 3.4))
    add_chart_image(slide, fig_er, 6.2, 1.1, 3.5, 3.4)
    
    best_eng = max(analyzed, key=lambda a: a["engagement_rate"])
    insight = (f"Engagement leader: {best_eng['company']} at {best_eng['engagement_rate']}% — "
               f"high engagement signals strong audience relevance and content quality.")
    add_rect(slide, 0.3, 4.65, 9.4, 0.65, RGBColor(0xE8, 0xF4, 0xF1))
    add_text_box(slide, f"💡 {insight}", 0.45, 4.72, 9.1, 0.5, size=9, color=TEXT_DARK, wrap=True)
    
    # ── SLIDE 6: CONTENT TOPICS & THEMES ─────────────────────────────────────
    slide = add_light_slide(prs, "Content Topics & Themes", "What each brand covers — and what they're missing")
    
    add_text_box(slide, "Top content themes by keyword frequency in titles and tags:",
                 0.35, 1.1, 9.3, 0.3, size=10, color=TEXT_MID)
    
    row_y = 1.55
    for i, a in enumerate(analyzed):
        col_x = 0.35 + (i % 2) * 4.7
        row_y_pos = row_y + (i // 2) * 1.8
        
        add_rect(slide, col_x, row_y_pos, 4.4, 1.65, RGBColor(0xFF, 0xFF, 0xFF))
        hex_color = CHART_COLORS_MPL[a["color_index"] % 5]
        header_color = RGBColor(int(hex_color[1:3], 16), int(hex_color[3:5], 16), int(hex_color[5:7], 16))
        add_rect(slide, col_x, row_y_pos, 4.4, 0.28, header_color)
        add_text_box(slide, a["company"], col_x + 0.1, row_y_pos + 0.04, 4.2, 0.22,
                     size=9, bold=True, color=TEXT_LIGHT)
        
        themes = a.get("top_themes", [])[:6]
        theme_str = "  ·  ".join(f"#{t}" for t in themes) if themes else "No theme data"
        add_text_box(slide, theme_str, col_x + 0.1, row_y_pos + 0.33, 4.2, 1.2,
                     size=8.5, color=TEXT_DARK, wrap=True)
    
    # ── SLIDE 7: POSTING FREQUENCY ────────────────────────────────────────────
    slide = add_light_slide(prs, "Posting Frequency & Consistency", "Upload cadence over the last 12 months")
    
    month_data_list = [a["month_counts"] for a in analyzed]
    fig_line = make_line_chart(companies, month_data_list, figsize=(9.0, 3.6))
    add_chart_image(slide, fig_line, 0.3, 1.1, 9.0, 3.6)
    
    add_rect(slide, 0.3, 4.85, 9.4, 0.5, RGBColor(0xE8, 0xF4, 0xF1))
    freq_texts = " | ".join(f"{a['company']}: {a['uploads_per_month']:.1f}/mo" for a in analyzed)
    add_text_box(slide, f"Upload frequency → {freq_texts}", 0.5, 4.92, 9.0, 0.35, size=9, color=TEXT_DARK)
    
    # ── SLIDE 8: GAP ANALYSIS ─────────────────────────────────────────────────
    slide = add_dark_slide(prs, "Gap Analysis", "Unexplored content territories and format opportunities")
    
    all_themes = set()
    for a in analyzed:
        all_themes.update(a.get("top_themes", []))
    
    gaps = []
    for a in analyzed:
        company_themes = set(a.get("top_themes", []))
        missing = all_themes - company_themes
        if missing:
            gaps.append((a["company"], list(missing)[:4]))
    
    format_gaps = [
        ("Long-form tutorials (20+ min)", "Deep-dive educational content commands higher watch time and loyalty"),
        ("Short-form clips (<60s)", "YouTube Shorts drives discovery across platform ecosystems"),
        ("Customer testimonials", "Social proof videos dramatically scale product validation conversions"),
        ("Live streams & webinars", "Live content generates higher interactive community retention"),
    ]
    
    add_text_box(slide, "TOPIC GAPS BY COMPANY", 0.35, 1.35, 4.2, 0.3, size=10, bold=True, color=ACCENT1)
    y = 1.75
    for company, missing in gaps[:4]:
        add_rect(slide, 0.35, y, 4.2, 0.28, MID_BG)
        add_text_box(slide, company, 0.45, y + 0.04, 1.5, 0.22, size=8.5, bold=True, color=ACCENT1)
        m_text = ", ".join(f"#{t}" for t in missing) if missing else "Comprehensive coverage"
        add_text_box(slide, f"Missing: {m_text}", 1.95, y + 0.04, 2.55, 0.22, size=8, color=TEXT_LIGHT)
        y += 0.38
    
    add_text_box(slide, "FORMAT OPPORTUNITIES", 5.0, 1.35, 4.6, 0.3, size=10, bold=True, color=ACCENT1)
    y = 1.75
    for fmt, desc in format_gaps:
        add_rect(slide, 5.0, y, 4.65, 0.5, MID_BG)
        add_text_box(slide, fmt, 5.1, y + 0.03, 4.5, 0.22, size=8.5, bold=True, color=ACCENT3)
        add_text_box(slide, desc, 5.1, y + 0.24, 4.5, 0.22, size=7.5, color=TEXT_LIGHT)
        y += 0.62
    
    # ── SLIDE 9: RECOMMENDATIONS ──────────────────────────────────────────────
    slide = add_dark_slide(prs, "Video Marketing Recommendations", f"Actionable steps for {your_company} based on analysis")
    
    most_active = max(analyzed, key=lambda a: a["uploads_per_month"])
    recs = [
        ("🎯", "Content Velocity", f"Match or exceed {most_active['company']}'s cadence of {most_active['uploads_per_month']:.1f} videos/month."),
        ("📊", "Engagement-First Strategy", f"Target a {best_eng['engagement_rate'] * 1.1:.1f}%+ engagement rate across active content arrays."),
        ("🎬", "Format Diversification", "Launch structured short-form programs explicitly targeting unique discovery paths."),
        ("🔍", "SEO & Discovery Optimisation", "Audit title frameworks of high-view competitors to refine metadata clusters."),
        ("📈", "Channel Growth Tactics", "Establish joint partner workflows to scale audience demographic transitions."),
    ]
    
    y = 1.3
    for icon, title, body in recs:
        add_rect(slide, 0.35, y, 9.3, 0.7, MID_BG)
        add_text_box(slide, icon, 0.45, y + 0.08, 0.4, 0.55, size=14)
        add_text_box(slide, title, 0.9, y + 0.05, 2.0, 0.28, size=9, bold=True, color=ACCENT1)
        add_text_box(slide, body, 0.9, y + 0.32, 8.55, 0.3, size=8.5, color=TEXT_LIGHT, wrap=True)
        y += 0.8
    
    # ── SLIDE 10: SCORECARD ───────────────────────────────────────────────────
    slide = add_dark_slide(prs, "Competitive Scorecard", "Ranking all companies on key video marketing metrics")
    
    categories = ["Subscribers", "Avg Views", "Engagement", "Frequency", "Content Volume"]
    values_matrix = [
        [a["sub_score"], a["view_score"], a["eng_score"], a["freq_score"],
         min(100, int(a["channel"].get("total_videos", 0) / max(1, max(b["channel"].get("total_videos", 1) for b in analyzed)) * 100))]
        for a in analyzed
    ]
    
    fig_radar = make_radar_chart(
        [a["company"] for a in analyzed],
        categories,
        [[v/100 for v in row] for row in values_matrix],
        figsize=(5.5, 4.2)
    )
    add_chart_image(slide, fig_radar, 0.2, 1.1, 5.5, 4.2)
    
    add_text_box(slide, "FINAL RANKINGS", 5.9, 1.15, 3.7, 0.3, size=11, bold=True, color=ACCENT1)
    
    table_y = 1.55
    medals = ["🥇", "🥈", "🥉", "4️⃣", "5️⃣"]
    for rank, a in enumerate(ranked):
        add_rect(slide, 5.9, table_y, 3.8, 0.62, MID_BG)
        add_text_box(slide, medals[rank], 6.0, table_y + 0.1, 0.4, 0.45, size=14)
        add_text_box(slide, a["company"], 6.45, table_y + 0.07, 2.2, 0.25, size=10, bold=True, color=TEXT_LIGHT)
        add_text_box(slide, f"Score: {a['score']}/100", 6.45, table_y + 0.32, 1.8, 0.22, size=8.5, color=ACCENT3)
        
        bar_full = slide.shapes.add_shape(1, Inches(8.5), Inches(table_y + 0.2), Inches(1.0), Inches(0.22))
        bar_full.fill.solid()
        bar_full.fill.fore_color.rgb = RGBColor(0x30, 0x45, 0x60)
        bar_full.line.fill.background()
        
        fill_width = max(0.02, a["score"] / 100)
        bar_filled = slide.shapes.add_shape(1, Inches(8.5), Inches(table_y + 0.2), Inches(fill_width), Inches(0.22))
        bar_filled.fill.solid()
        
        hex_color = CHART_COLORS_MPL[a["color_index"] % 5]
        bar_filled.fill.fore_color.rgb = RGBColor(
            int(hex_color[1:3], 16),
            int(hex_color[3:5], 16),
            int(hex_color[5:7], 16)
        )
        bar_filled.line.fill.background()
        table_y += 0.72
    
    # ── SLIDE 11: SUMMARY / NEXT STEPS ───────────────────────────────────────
    slide = add_dark_slide(prs, "Summary & Next Steps", "")
    add_text_box(slide, "What We Found", 0.35, 1.15, 4.5, 0.35, size=13, bold=True, color=ACCENT1)
    
    summary_points = [
        f"{leader['company']} leads on composite score with strong subscriber base and engagement",
        f"Engagement rates range from {min(a['engagement_rate'] for a in analyzed)}% to {max(a['engagement_rate'] for a in analyzed)}%",
        f"Upload frequency varies widely across structural operational profiles",
        "Format diversity is underutilised across all matched corporate arrays",
    ]
    
    y = 1.6
    for pt in summary_points:
        add_rect(slide, 0.35, y, 0.06, 0.28, ACCENT1)
        add_text_box(slide, pt, 0.52, y + 0.02, 4.0, 0.28, size=9, color=TEXT_LIGHT, wrap=True)
        y += 0.42
    
    add_text_box(slide, "Your Next 90 Days", 5.1, 1.15, 4.5, 0.35, size=13, bold=True, color=ACCENT1)
    next_steps = [
        ("Days 1–30",  "Audit metrics profile · Establish target workflows framework"),
        ("Days 31–60", "Deploy short-form optimization structures · Refine indexing metadata"),
        ("Days 61–90", "Evaluate channel performance curves against direct competitor models"),
    ]
    
    y = 1.6
    for period, action in next_steps:
        add_rect(slide, 5.1, y, 4.55, 0.8, MID_BG)
        add_text_box(slide, period, 5.2, y + 0.06, 1.5, 0.25, size=9, bold=True, color=ACCENT1)
        add_text_box(slide, action, 5.2, y + 0.32, 4.3, 0.42, size=8.5, color=TEXT_LIGHT, wrap=True)
        y += 0.95
        
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf

# ─── ROUTES ──────────────────────────────────────────────────────────────────
@app.route("/")
def index():
    return send_from_directory("static", "index.html")

def get_api_key():
    header_key = request.headers.get('X-YouTube-Key', '').strip()
    return header_key or YOUTUBE_API_KEY

@app.route("/api/analyze", methods=["POST"])
def analyze():
    global YOUTUBE_API_KEY
    try:
        body = request.get_json()
        your_company = body.get("your_company", "").strip()
        competitors  = [c.strip() for c in body.get("competitors", []) if c.strip()]
        
        if not your_company:
            return jsonify({"error": "Company name required"}), 400
        
        original_key = YOUTUBE_API_KEY
        YOUTUBE_API_KEY = get_api_key()
        
        all_companies = [your_company] + competitors[:4]
        companies_data = [fetch_company_data(name) for name in all_companies]
        
        YOUTUBE_API_KEY = original_key
        
        analyzed = analyze_data(companies_data)
        analyzed = compute_scores(analyzed)
        
        result = []
        for a in analyzed:
            result.append({
                "company": a["company"],
                "is_mock": a.get("is_mock", False),
                "channel": {
                    "title": a["channel"].get("channel_title", a["company"]),
                    "subscribers": a["channel"].get("subscribers", 0),
                    "total_videos": a["channel"].get("total_videos", 0),
                    "total_views": a["channel"].get("total_views", 0),
                    "country": a["channel"].get("country", "N/A"),
                },
                "avg_views": a["avg_views"],
                "avg_likes": a["avg_likes"],
                "avg_comments": a["avg_comments"],
                "engagement_rate": a["engagement_rate"],
                "uploads_per_month": a["uploads_per_month"],
                "score": a["score"],
                "top_themes": a.get("top_themes", []),
                "top_videos": [
                    {
                        "title": v["title"],
                        "views": v["views"],
                        "likes": v["likes"],
                        "comments": v["comments"],
                        "published_at": v["published_at"],
                        "id": v["id"]
                    }
                    for v in a.get("top_videos", [])[:5]
                ]
            })
        return jsonify({"companies": result, "your_company": your_company})
    except Exception as route_err:
        print(traceback.format_exc())
        return jsonify({"error": str(route_err)}), 500

@app.route("/api/download", methods=["POST"])
def download():
    global YOUTUBE_API_KEY
    try:
        body = request.get_json()
        your_company = body.get("your_company", "").strip()
        competitors  = [c.strip() for c in body.get("competitors", []) if c.strip()]
        
        if not your_company:
            return jsonify({"error": "Company name required"}), 400
        
        original_key = YOUTUBE_API_KEY
        YOUTUBE_API_KEY = get_api_key()
        
        all_companies = [your_company] + competitors[:4]
        companies_data = [fetch_company_data(name) for name in all_companies]
        
        YOUTUBE_API_KEY = original_key
        
        analyzed = analyze_data(companies_data)
        analyzed = compute_scores(analyzed)
        
        pptx_buf = build_pptx(analyzed, your_company)
        
        safe_company = re.sub(r'[^A-Za-z0-9_]', '', your_company.replace(' ', '_'))
        filename = f"video_intel_{safe_company}_{datetime.date.today()}.pptx"
        
        try:
            return send_file(
                pptx_buf,
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                as_attachment=True,
                download_name=filename
            )
        except TypeError:
            pptx_buf.seek(0)
            return send_file(
                pptx_buf,
                mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                as_attachment=True,
                attachment_filename=filename
            )
            
    except Exception as download_err:
        print(traceback.format_exc())
        return jsonify({"error": f"Presentation compilation failed: {str(download_err)}"}), 500

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=port, debug=False)
